"""
Generic Text Extractor
======================

Dispatches to the appropriate extractor based on file extension.
Every connector eventually feeds through this layer.

Supported:
  .pdf  -> PDFExtractor (pdfplumber)
  .txt  -> plain text read
  .md   -> plain text read (Markdown is just text)
  .docx -> DOCXExtractor (python-docx)
  .xlsx -> ExcelExtractor (openpyxl + pandas)
  .xls  -> Legacy ExcelExtractor (xlrd + pandas)
  .csv  -> CSVExtractor (pandas)
  .pptx -> PPTXExtractor (python-pptx)
  .html -> HTML extractor (html.parser)
  .json -> JSON text extraction
  .xml  -> XML text extraction (xml.etree)

Unsupported extensions return None (caller decides fallback).
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Optional


class TextExtractor:
    """
    Strategy-pattern text extractor.

    Usage:
        extractor = TextExtractor()
        text = extractor.extract(Path("document.pdf"))
    """

    def extract(self, file_path: Path) -> Optional[str]:
        """
        Extract text from a file.

        Returns extracted text, or None if the file type is unsupported.
        Raises FileNotFoundError if file does not exist.
        """
        if not file_path.exists():
            raise FileNotFoundError(file_path)

        suffix = file_path.suffix.lower()

        if suffix == ".pdf":
            return self._extract_pdf(file_path)
        elif suffix in (".txt", ".md", ".markdown", ".rst"):
            return self._extract_text(file_path)
        elif suffix == ".docx":
            return self._extract_docx(file_path)
        elif suffix == ".xlsx":
            return self._extract_excel(file_path)
        elif suffix == ".xls":
            return self._extract_xls(file_path)
        elif suffix == ".csv":
            return self._extract_csv(file_path)
        elif suffix == ".pptx":
            return self._extract_pptx(file_path)
        elif suffix in (".html", ".htm"):
            return self._extract_html(file_path)
        elif suffix == ".json":
            return self._extract_json(file_path)
        elif suffix == ".xml":
            return self._extract_xml(file_path)
        else:
            return None

    @staticmethod
    def supported_extensions() -> set[str]:
        """Return the set of supported file extensions."""
        return {".pdf", ".txt", ".md", ".markdown", ".rst",
                ".docx", ".xlsx", ".xls", ".csv",
                ".pptx", ".html", ".htm", ".json", ".xml"}

    # ---- Private extractors ----

    @staticmethod
    def _extract_pdf(file_path: Path) -> str:
        import pdfplumber
        pages: list[str] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(text)
        return "\n".join(pages)

    @staticmethod
    def _extract_text(file_path: Path) -> str:
        # Try UTF-8 first, fall back to system default
        for encoding in ("utf-8", "utf-8-sig", "latin-1"):
            try:
                return file_path.read_text(encoding=encoding)
            except (UnicodeDecodeError, ValueError):
                continue
        # Last resort: binary read with error replacement
        return file_path.read_bytes().decode("utf-8", errors="replace")

    @staticmethod
    def _extract_docx(file_path: Path) -> str:
        from docx import Document
        doc = Document(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract table text
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        return "\n".join(paragraphs)

    @staticmethod
    def _clean_excel_text(text: str) -> str:
        """Remove NaN and Unnamed artifacts from extracted spreadsheet text.

        Preserves cell content and structure while removing noise.
        """
        text = re.sub(r"Unnamed:\s*\d+", "", text)
        text = re.sub(r"\bNaN\b", "", text)
        return text

    @staticmethod
    def _extract_excel(file_path: Path) -> str:
        """Extract from .xlsx using openpyxl."""
        import pandas as pd
        frames = pd.read_excel(str(file_path), sheet_name=None, engine="openpyxl")
        parts: list[str] = []
        for sheet_name, df in frames.items():
            df_clean = df.dropna(how="all")
            if df_clean.empty:
                continue
            text = df_clean.to_string(index=False)
            text = TextExtractor._clean_excel_text(text)
            parts.append(f"--- Sheet: {sheet_name} ---")
            parts.append(text)
        return "\n".join(parts)

    @staticmethod
    def _extract_xls(file_path: Path) -> str:
        """Extract from legacy .xls using xlrd."""
        import pandas as pd
        frames = pd.read_excel(str(file_path), sheet_name=None, engine="xlrd")
        parts: list[str] = []
        for sheet_name, df in frames.items():
            df_clean = df.dropna(how="all")
            if df_clean.empty:
                continue
            text = df_clean.to_string(index=False)
            text = TextExtractor._clean_excel_text(text)
            parts.append(f"--- Sheet: {sheet_name} ---")
            parts.append(text)
        return "\n".join(parts)

    @staticmethod
    def _extract_csv(file_path: Path) -> str:
        import pandas as pd
        df = pd.read_csv(str(file_path))
        return df.to_string(index=False)

    @staticmethod
    def _extract_pptx(file_path: Path) -> str:
        """Extract text from PowerPoint (.pptx) files."""
        from pptx import Presentation
        prs = Presentation(str(file_path))
        parts: list[str] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)
            if slide_texts:
                parts.append(f"--- Slide {slide_num} ---")
                parts.append("\n".join(slide_texts))
        return "\n".join(parts)

    @staticmethod
    def _extract_html(file_path: Path) -> str:
        """Extract text from HTML files using html.parser."""
        from html.parser import HTMLParser
        import io

        class HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self._result: list[str] = []
                self._skip = False

            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style", "noscript"):
                    self._skip = True

            def handle_endtag(self, tag):
                if tag in ("script", "style", "noscript"):
                    self._skip = False

            def handle_data(self, data):
                if not self._skip:
                    text = data.strip()
                    if text:
                        self._result.append(text)

            def get_text(self) -> str:
                return " ".join(self._result)

        content = TextExtractor._extract_text(file_path)
        extractor = HTMLTextExtractor()
        extractor.feed(content)
        return extractor.get_text()

    @staticmethod
    def _extract_json(file_path: Path) -> str:
        """Extract readable text from JSON files."""
        import json
        content = TextExtractor._extract_text(file_path)
        try:
            data = json.loads(content)
            return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            return content

    @staticmethod
    def _extract_xml(file_path: Path) -> str:
        """Extract text content from XML files."""
        import xml.etree.ElementTree as ET
        content = TextExtractor._extract_text(file_path)
        try:
            root = ET.fromstring(content)
            texts: list[str] = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    texts.append(elem.tail.strip())
            return " ".join(texts)
        except ET.ParseError:
            return content
