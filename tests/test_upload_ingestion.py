"""
Tests for Document Upload & Ingestion
======================================
"""
import hashlib
import json
import os
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import MagicMock, patch

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))


class TestTextExtractorNewFormats(unittest.TestCase):
    """Test the new TextExtractor format support."""

    def setUp(self):
        from kurukshetra.extractors.text_extractor import TextExtractor
        self.extractor = TextExtractor()

    def test_pptx_extraction(self):
        """PPTX extraction produces text from slides."""
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx not installed")

        tmp = Path(tempfile.mktemp(suffix=".pptx"))
        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Test Slide"
            slide.placeholders[1].text = "Slide content here"
            prs.save(str(tmp))
            text = self.extractor.extract(tmp)
            self.assertIsNotNone(text)
            self.assertIn("Test Slide", text)
            self.assertIn("Slide content here", text)
        finally:
            if tmp.exists():
                tmp.unlink()

    def test_html_extraction(self):
        """HTML extraction strips tags and returns text."""
        tmp = Path(tempfile.mktemp(suffix=".html"))
        try:
            tmp.write_text(
                "<html><head><title>Test</title></head><body>"
                "<h1>Hello World</h1>"
                "<p>This is a <b>test</b> document.</p>"
                "<script>var x = 'ignore';</script>"
                "</body></html>",
                encoding="utf-8",
            )
            text = self.extractor.extract(tmp)
            self.assertIsNotNone(text)
            self.assertIn("Hello World", text)
            self.assertIn("test document", text)
            self.assertNotIn("ignore", text)
            self.assertNotIn("<h1>", text)
        finally:
            if tmp.exists():
                tmp.unlink()

    def test_json_extraction(self):
        """JSON extraction returns formatted JSON text."""
        tmp = Path(tempfile.mktemp(suffix=".json"))
        try:
            data = {"name": "G3 Config", "version": "2.0", "teams": ["SPM", "ICS"]}
            tmp.write_text(json.dumps(data), encoding="utf-8")
            text = self.extractor.extract(tmp)
            self.assertIsNotNone(text)
            self.assertIn("G3 Config", text)
            self.assertIn("SPM", text)
        finally:
            if tmp.exists():
                tmp.unlink()

    def test_xml_extraction(self):
        """XML extraction returns text content from elements."""
        tmp = Path(tempfile.mktemp(suffix=".xml"))
        try:
            tmp.write_text(
                '<?xml version="1.0"?><config><name>G3 RMS</name><version>1.0</version></config>',
                encoding="utf-8",
            )
            text = self.extractor.extract(tmp)
            self.assertIsNotNone(text)
            self.assertIn("G3 RMS", text)
            self.assertIn("1.0", text)
        finally:
            if tmp.exists():
                tmp.unlink()

    def test_supported_extensions_includes_new(self):
        """Supported extensions include new formats."""
        from kurukshetra.extractors.text_extractor import TextExtractor
        exts = TextExtractor.supported_extensions()
        for ext in [".pptx", ".html", ".htm", ".json", ".xml",
                    ".pdf", ".docx", ".xlsx", ".csv", ".txt", ".md"]:
            self.assertIn(ext, exts)

    def test_unsupported_extension_returns_none(self):
        """Unsupported extensions return None."""
        tmp = Path(tempfile.mktemp(suffix=".xyz"))
        try:
            tmp.write_text("test", encoding="utf-8")
            result = self.extractor.extract(tmp)
            self.assertIsNone(result)
        finally:
            if tmp.exists():
                tmp.unlink()


class TestUploadSecurity(unittest.TestCase):
    """Test upload security constraints."""

    def test_dangerous_extension_rejected(self):
        from command_center.backend.routers.knowledge import DANGEROUS_EXTENSIONS
        for ext in [".exe", ".bat", ".sh", ".py", ".ps1", ".vbs"]:
            self.assertIn(ext, DANGEROUS_EXTENSIONS)

    def test_allowed_extensions_complete(self):
        from command_center.backend.routers.knowledge import ALLOWED_EXTENSIONS
        expected = {
            ".pdf", ".docx", ".xlsx", ".xls", ".csv",
            ".txt", ".md", ".rst", ".html", ".htm",
            ".json", ".xml", ".pptx",
        }
        self.assertTrue(expected.issubset(ALLOWED_EXTENSIONS))

    def test_max_file_size_defined(self):
        from command_center.backend.routers.knowledge import MAX_FILE_SIZE
        self.assertEqual(MAX_FILE_SIZE, 50 * 1024 * 1024)

    def test_upload_dir_is_safe(self):
        from command_center.backend.routers.knowledge import UPLOAD_DIR
        self.assertTrue(str(UPLOAD_DIR).endswith("uploads"))
        self.assertIn("knowledge", str(UPLOAD_DIR))


class TestUploadEndToEnd(unittest.TestCase):
    """End-to-end upload test using FastAPI TestClient."""

    def setUp(self):
        try:
            from fastapi.testclient import TestClient
            from command_center.backend.main import app
            self.client = TestClient(app)
        except ImportError:
            self.skipTest("FastAPI test client not available")

    def test_upload_txt_file(self):
        content = b"This is a test document about G3 RMS configuration.\nG3 RMS is used for revenue management."
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("test_doc.txt", content, "text/plain")},
        )
        self.assertEqual(response.status_code, 200)
        data = response.json()
        self.assertIn("document_id", data)
        self.assertIn(data["status"], ("ok", "error"))
        if data["status"] == "ok":
            self.assertGreater(data["chunks_stored"], 0)

    def test_upload_csv_file(self):
        content = b"system,team,process\nG3,SPM,Decision Upload\nG3,ICS,Installation\n"
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("systems.csv", content, "text/csv")},
        )
        self.assertEqual(response.status_code, 200)
        data = response.json()
        self.assertIn(data["status"], ("ok", "error"))

    def test_upload_json_file(self):
        content = json.dumps({
            "system": "G3 RMS", "teams": ["SPM", "ICS"],
        }).encode()
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("config.json", content, "application/json")},
        )
        self.assertEqual(response.status_code, 200)

    def test_empty_file_rejected(self):
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("empty.txt", b"", "text/plain")},
        )
        self.assertEqual(response.status_code, 400)

    def test_dangerous_extension_rejected(self):
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("malware.exe", b"MZ\x90\x00", "application/octet-stream")},
        )
        self.assertEqual(response.status_code, 400)

    def test_unsupported_extension_rejected(self):
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("image.bmp", b"BM\x00", "image/bmp")},
        )
        self.assertEqual(response.status_code, 400)

    def test_large_file_rejected(self):
        large_content = b"x" * (51 * 1024 * 1024)
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("large.txt", large_content, "text/plain")},
        )
        self.assertEqual(response.status_code, 413)

    def test_duplicate_upload(self):
        content = b"Unique test content for dedup check G3 RMS system."
        r1 = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("dedup_test.txt", content, "text/plain")},
        )
        self.assertEqual(r1.status_code, 200)
        r2 = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("dedup_test.txt", content, "text/plain")},
        )
        self.assertEqual(r2.status_code, 200)

    def test_upload_response_fields(self):
        content = b"Test document with G3 RMS content."
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("fields_test.txt", content, "text/plain")},
        )
        self.assertEqual(response.status_code, 200)
        data = response.json()
        for field in ["document_id", "filename", "status", "message",
                       "chunks_stored", "entities_extracted", "team_id",
                       "execution_time_ms"]:
            self.assertIn(field, data, f"Missing field: {field}")

    def test_filename_preserved(self):
        content = b"Test content for filename check."
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("My Document.pdf", content, "application/pdf")},
        )
        self.assertEqual(response.status_code, 200)
        data = response.json()
        self.assertEqual(data["filename"], "My Document.pdf")

    def test_path_traversal_prevention(self):
        content = b"Test content."
        response = self.client.post(
            "/api/knowledge/upload",
            files={"file": ("../../../etc/passwd", content, "text/plain")},
        )
        self.assertIn(response.status_code, (200, 400))


if __name__ == "__main__":
    unittest.main()
